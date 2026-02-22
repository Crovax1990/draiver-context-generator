"""
Draiver Context Generator – PPTX Renderer
Generates PowerPoint files using python-pptx templates.
"""
import logging
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from src.rag_engine import SlideContent

logger = logging.getLogger(__name__)

# Constants for layout indices based on standard templates
LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_SECTION_HEADER = 2
LAYOUT_TWO_CONTENT = 3  # Often used for text + image
LAYOUT_TITLE_ONLY = 4
LAYOUT_BLANK = 5

class PPTXRenderer:
    def __init__(self, template_path: Optional[Path] = None):
        if template_path and template_path.exists():
            logger.info("Using template from %s", template_path)
            self.prs = Presentation(str(template_path))
        else:
            logger.info("No template found, creating new presentation from scratch")
            self.prs = Presentation()
            # Set default slide size to 16:9
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title: str, subtitle: str):
        layout = self.prs.slide_layouts[LAYOUT_TITLE]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return slide

    def add_section_header(self, title: str):
        layout = self.prs.slide_layouts[LAYOUT_SECTION_HEADER]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        return slide

    def add_content_slide(
        self, 
        slide_data: SlideContent, 
        image_path: Optional[str] = None
    ):
        """Adds a slide with bullet points and optionally an image."""
        # Choose layout based on image presence
        # If image exists, use Two Content (idx 3), else Title and Content (idx 1)
        layout_idx = LAYOUT_TWO_CONTENT if image_path else LAYOUT_TITLE_CONTENT
        
        # Fallback if layout doesn't exist in template
        if layout_idx >= len(self.prs.slide_layouts):
            layout_idx = LAYOUT_TITLE_CONTENT

        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        
        # Title
        slide.shapes.title.text = slide_data.titolo_slide
        
        # Bullet points
        # Placeholder 1 is typically the body text frame
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default text
        
        for i, bullet in enumerate(slide_data.bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            # Basic styling
            p.font.size = Pt(20)
            p.font.name = "Calibri"

        # Image
        if image_path:
            # For Two Content layout, placeholder 2 is usually the second content box
            # If not, we manually add the picture to the right side
            try:
                if len(slide.placeholders) > 2:
                    placeholder = slide.placeholders[2]
                    # Get dimensions of placeholder
                    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
                    # Remove placeholder before adding image to avoid overlap
                    # Note: python-pptx makes it hard to 'fill' a placeholder with an image directly
                    # so we just use its coordinates
                    # slide.shapes.add_picture(image_path, left, top, width=width)
                    
                    # Manually position on the right half
                    slide.shapes.add_picture(
                        image_path, 
                        Inches(7), Inches(1.5), 
                        width=Inches(5.5)
                    )
                else:
                    # Fallback manual positioning
                    slide.shapes.add_picture(
                        image_path, 
                        Inches(7.5), Inches(1.5), 
                        width=Inches(5)
                    )
            except Exception as e:
                logger.error("Error adding image %s: %s", image_path, e)

        # Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data.note_relatore
        if slide_data.source_doc_names:
            text_frame.add_paragraph().text = f"Fonti: {', '.join(slide_data.source_doc_names)}"

        return slide

    def add_final_slide(self, title: str, content: List[str]):
        layout = self.prs.slide_layouts[LAYOUT_TITLE_CONTENT]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in content:
            p = tf.add_paragraph()
            p.text = bullet
        return slide

    def save(self, output_path: Path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("Presentation saved to %s", output_path)
