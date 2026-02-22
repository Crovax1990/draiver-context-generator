from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_professional_template(output_path):
    prs = Presentation()

    # Define some professional colors
    UNIV_BLUE = RGBColor(0, 51, 102) # Classic Academic Blue
    ACCENT_GRAY = RGBColor(128, 128, 128)

    # 1. Title Slide Layout (Index 0)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Università degli Studi - Corso di Ostetricia"
    subtitle.text = "Presentazione della Lezione\nSottotitolo della Sessione"
    
    # Simple styling for title slide
    title.text_frame.paragraphs[0].font.color.rgb = UNIV_BLUE
    title.text_frame.paragraphs[0].font.bold = True

    # 2. Section Header Layout (Index 2)
    section_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(section_layout)
    title = slide.shapes.title
    title.text = "Sezione I: Introduzione"
    title.text_frame.paragraphs[0].font.color.rgb = UNIV_BLUE

    # 3. Title and Content Layout (Index 1)
    content_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Titolo della Slide di Contenuto"
    tf = body.text_frame
    tf.text = "Punto principale della discussione"
    
    p = tf.add_paragraph()
    p.text = "Sotto-argomento dettagliato"
    p.level = 1

    # 4. Two Content Layout (Index 3)
    two_content_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(two_content_layout)
    title = slide.shapes.title
    title.text = "Confronto o Immagine e Testo"
    
    left_placeholder = slide.placeholders[1]
    right_placeholder = slide.placeholders[2]
    
    left_placeholder.text = "Descrizione testuale a sinistra"
    right_placeholder.text = "[Spazio per Immagine o Grafico a destra]"

    # Save the template
    prs.save(output_path)
    print(f"✅ Professional template created at {output_path}")

if __name__ == "__main__":
    create_professional_template("template.pptx")
